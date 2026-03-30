import requests
import os
from pptx import Presentation
import openpyxl

def create_test_files():
    # Create dummy PPTX
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Test PowerPoint Support"
    prs.save("test.pptx")
    
    # Create dummy XLSX
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Test Excel Support"
    ws["A2"] = "Data Row 1"
    wb.save("test.xlsx")

def test_upload(filename):
    url = "http://127.0.0.1:7860/api/upload"
    try:
        with open(filename, "rb") as f:
            files = {'file': (filename, f)}
            response = requests.post(url, files=files)
            print(f"File: {filename} | Status: {response.status_code}")
            print(f"Response: {response.json()}")
    except Exception as e:
        print(f"Error uploading {filename}: {e}")

if __name__ == "__main__":
    create_test_files()
    test_upload("test.pptx")
    test_upload("test.xlsx")
    
    # Cleanup
    if os.path.exists("test.pptx"): os.remove("test.pptx")
    if os.path.exists("test.xlsx"): os.remove("test.xlsx")

<<<<<<< HEAD
from PyPDF2 import PdfReader
from pptx import Presentation
import openpyxl


def load_pdf(uploaded_file):
    pdf_reader = PdfReader(uploaded_file)
    text = ""

    for page in pdf_reader.pages:
        extracted = page.extract_text()
        if extracted:
            text += extracted + "\n"

    return text


def load_pptx(uploaded_file):
    prs = Presentation(uploaded_file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def load_xlsx(uploaded_file):
    wb = openpyxl.load_workbook(uploaded_file, data_only=True)
    text = []
    for sheet in wb.sheetnames:
        text.append(f"--- Sheet: {sheet} ---")
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
            if row_text.strip():
                text.append(row_text)
    return "\n".join(text)


def load_multiple_pdfs(uploaded_files):
    all_text = ""

    for file in uploaded_files:
        text = load_pdf(file)
        all_text += text + "\n"

    return all_text
=======
"""
document_loader.py — Document Processing Pipeline

WHY THIS MODULE EXISTS:
    RAG systems need clean, chunked text with metadata to work well.
    This module handles the entire ingestion pipeline:
    1. Extract raw text from PDF/DOCX files
    2. Detect the document type (notes, question paper, lab manual)
    3. Attach metadata (source, type, page) to each chunk
    4. Split text into optimally-sized chunks for embedding

WHY THESE LIBRARIES:
    - PyPDF2: Lightweight, pure-Python PDF reader. No system dependencies.
    - python-docx: Standard library for Word document parsing.
    - RecursiveCharacterTextSplitter: LangChain's best general-purpose splitter.
      It tries to split on paragraphs → sentences → words, preserving meaning.
"""

import os
import re
import logging
from pathlib import Path

from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

from config import CHUNK_SIZE, CHUNK_OVERLAP, ALLOWED_EXTENSIONS

logger = logging.getLogger(__name__)


# ── Document Type Detection ──────────────────────────────────────────────────
# WHY: Knowing the document type lets the LLM tailor its response.
#      A question paper needs different handling than lecture notes.

# Keyword patterns for classification — easily extensible
_DOC_TYPE_PATTERNS: dict[str, list[str]] = {
    "question_paper": [
        r"question\s*paper",
        r"examination",
        r"marks?\s*:",
        r"attempt\s+(any|all)",
        r"max\.?\s*marks",
        r"time\s*:\s*\d",
        r"answer\s+(any|all)",
        r"instructions\s+to\s+candidates",
    ],
    "lab_manual": [
        r"experiment\s*no",
        r"lab\s*(manual|report)",
        r"apparatus",
        r"procedure\s*:",
        r"observation\s*table",
        r"aim\s*:",
        r"theory\s*:",
        r"result\s*:",
    ],
    "notes": [
        r"chapter\s*\d",
        r"unit\s*\d",
        r"module\s*\d",
        r"lecture\s*notes",
        r"summary\s*:",
        r"key\s*points",
        r"introduction\s*:",
        r"definition\s*:",
    ],
}


def detect_document_type(text: str) -> str:
    """
    Classify a document based on keyword pattern matching.

    WHY keyword-based instead of ML: Simpler, deterministic, zero latency,
    and works well for structured academic documents. Can be upgraded to
    an ML classifier in Track B.

    Args:
        text: First ~2000 characters of the document for classification.

    Returns:
        One of: 'question_paper', 'lab_manual', 'notes', 'general'.
    """
    sample = text[:2000].lower()
    scores: dict[str, int] = {}

    for doc_type, patterns in _DOC_TYPE_PATTERNS.items():
        score = sum(
            1 for pattern in patterns if re.search(pattern, sample, re.IGNORECASE)
        )
        scores[doc_type] = score

    best_type = max(scores, key=scores.get)  # type: ignore[arg-type]
    return best_type if scores[best_type] >= 2 else "general"


# ── Text Extraction ──────────────────────────────────────────────────────────


def _extract_pdf(file_path: str) -> list[dict[str, str | int]]:
    """
    Extract text page-by-page from a PDF.

    WHY page-by-page: Preserves page numbers for source attribution.
    Each page becomes a separate unit before chunking.

    Returns:
        List of dicts with 'text' and 'page' keys.
    """
    reader = PdfReader(file_path)
    pages: list[dict[str, str | int]] = []

    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            pages.append({"text": text.strip(), "page": i + 1})

    if not pages:
        logger.warning("No text extracted from PDF: %s", file_path)

    return pages


def _extract_docx(file_path: str) -> list[dict[str, str | int]]:
    """
    Extract text from a DOCX file.

    WHY paragraph-based: DOCX files don't have page numbers in the same way.
    We treat the entire document as page 1 for simplicity.
    Track B improvement: Use section/heading detection for finer granularity.

    Returns:
        List of dicts with 'text' and 'page' keys.
    """
    doc = DocxDocument(file_path)
    full_text = "\n".join(
        para.text.strip() for para in doc.paragraphs if para.text.strip()
    )

    if not full_text:
        logger.warning("No text extracted from DOCX: %s", file_path)
        return []

    return [{"text": full_text, "page": 1}]


# ── Chunking & Metadata ─────────────────────────────────────────────────────


def load_and_chunk(file_path: str) -> list[Document]:
    """
    Complete ingestion pipeline: extract → detect type → chunk → add metadata.

    WHY RecursiveCharacterTextSplitter:
        It tries to split on natural boundaries (\\n\\n → \\n → . → space).
        This preserves semantic meaning within chunks better than fixed-size
        splitting. The overlap ensures context isn't lost at boundaries.

    Args:
        file_path: Absolute path to a PDF or DOCX file.

    Returns:
        List of LangChain Document objects, each with metadata:
        - source: filename
        - doc_type: detected document type
        - page: page number (from extraction)
        - chunk_index: position within the chunk sequence

    Raises:
        ValueError: If file extension is not supported.
        FileNotFoundError: If file does not exist.
    """
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = path.suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type: {ext}. Allowed: {ALLOWED_EXTENSIONS}"
        )

    # Step 1: Extract raw text with page info
    if ext == ".pdf":
        pages = _extract_pdf(file_path)
    else:
        pages = _extract_docx(file_path)

    if not pages:
        return []

    # Step 2: Detect document type from first page's text
    full_text = " ".join(p["text"] for p in pages)
    doc_type = detect_document_type(full_text)
    logger.info("Detected document type for '%s': %s", path.name, doc_type)

    # Step 3: Create LangChain Documents with metadata
    raw_documents: list[Document] = []
    for page_data in pages:
        raw_documents.append(
            Document(
                page_content=str(page_data["text"]),
                metadata={
                    "source": path.name,
                    "doc_type": doc_type,
                    "page": page_data["page"],
                },
            )
        )

    # Step 4: Split into optimized chunks
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        length_function=len,
        separators=["\n\n", "\n", ". ", " ", ""],
    )

    chunks = splitter.split_documents(raw_documents)

    # Step 5: Add chunk index to metadata for traceability
    for i, chunk in enumerate(chunks):
        chunk.metadata["chunk_index"] = i

    logger.info(
        "Processed '%s': %d pages → %d chunks (type: %s)",
        path.name,
        len(pages),
        len(chunks),
        doc_type,
    )

    return chunks
>>>>>>> 5caf138cd4bf724bc60bc7be180102528e7f7762
